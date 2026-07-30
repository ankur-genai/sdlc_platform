import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck(filename, title_text, slides_data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    BG_DARK = RGBColor(18, 22, 33)       # Deep Charcoal / Navy #121621
    CARD_BG = RGBColor(27, 33, 48)       # Soft Dark Card #1B2130
    CARD_BORDER = RGBColor(45, 55, 78)   # Card Border #2D374E
    ACCENT_YELLOW = RGBColor(255, 230, 0) # EY Gold #FFE600
    ACCENT_BLUE = RGBColor(0, 229, 255)   # Vivid Cyan #00E5FF
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(160, 174, 192)

    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # 1. Background fill
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_DARK
        bg_shape.line.fill.background()

        # 2. Header Banner / Title
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        
        # Category Tag
        p_tag = tf.paragraphs[0]
        p_tag.text = slide_info.get("category", "EXECUTIVE DECK").upper()
        p_tag.font.name = "Calibri"
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = ACCENT_YELLOW
        p_tag.space_after = Pt(4)

        # Slide Title
        p_title = tf.add_paragraph()
        p_title.text = slide_info["title"]
        p_title.font.name = "Arial"
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.space_after = Pt(2)

        # Subtitle
        if slide_info.get("subtitle"):
            p_sub = tf.add_paragraph()
            p_sub.text = slide_info["subtitle"]
            p_sub.font.name = "Calibri"
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = TEXT_MUTED

        # 3. Render Cards / Content Grid
        cards = slide_info.get("cards", [])
        num_cards = len(cards)
        
        if num_cards > 0:
            margin_left = Inches(0.8)
            available_width = Inches(11.733)
            gap = Inches(0.3)
            card_width = (available_width - (gap * (num_cards - 1))) / num_cards
            top_pos = Inches(1.7)
            card_height = Inches(5.2)

            for i, card in enumerate(cards):
                left_pos = margin_left + i * (card_width + gap)

                # Card Background Container
                card_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, card_width, card_height
                )
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = CARD_BG
                card_shape.line.color.rgb = CARD_BORDER
                card_shape.line.width = Pt(1.5)

                # Card Content Text
                tb = slide.shapes.add_textbox(
                    left_pos + Inches(0.2), top_pos + Inches(0.2), card_width - Inches(0.4), card_height - Inches(0.4)
                )
                tf_card = tb.text_frame
                tf_card.word_wrap = True

                # Card Header / Metric
                p_c_head = tf_card.paragraphs[0]
                p_c_head.text = card.get("title", "")
                p_c_head.font.name = "Arial"
                p_c_head.font.size = Pt(16)
                p_c_head.font.bold = True
                p_c_head.font.color.rgb = ACCENT_BLUE if i % 2 == 1 else ACCENT_YELLOW
                p_c_head.space_after = Pt(10)

                # Big Metric / KPI callout if present
                if card.get("metric"):
                    p_m = tf_card.add_paragraph()
                    p_m.text = card["metric"]
                    p_m.font.name = "Arial"
                    p_m.font.size = Pt(28)
                    p_m.font.bold = True
                    p_m.font.color.rgb = ACCENT_YELLOW
                    p_m.space_after = Pt(8)

                # Bullets / Details
                for bullet in card.get("bullets", []):
                    p_b = tf_card.add_paragraph()
                    p_b.text = f"• {bullet}"
                    p_b.font.name = "Calibri"
                    p_b.font.size = Pt(12)
                    p_b.font.color.rgb = TEXT_WHITE
                    p_b.space_after = Pt(6)

        # 4. Speaker Notes
        notes_slide = slide.notes_slide
        tf_notes = notes_slide.notes_text_frame
        tf_notes.text = slide_info.get("speaker_notes", "")

    # Save output deck
    output_path = os.path.join(os.getcwd(), filename)
    prs.save(output_path)
    return output_path

def generate_all_decks():
    # ── DECK 1: Executive Overview ───────────────────────────────────────────
    exec_slides = [
        {
            "category": "EXECUTIVE SUMMARY",
            "title": "Autonomous SDLC Platform: Transforming Enterprise Engineering",
            "subtitle": "Accelerating software delivery with 15 specialized AI agents and human-in-the-loop governance.",
            "cards": [
                {
                    "title": "Core Mission",
                    "metric": "10x Speed",
                    "bullets": [
                        "Automate repetitive SDLC tasks from intake to deployment.",
                        "Enforce institutional security & quality standards automatically.",
                        "Reduce delivery lead time from weeks to minutes."
                    ]
                },
                {
                    "title": "AI Capability",
                    "metric": "15 Specialized Agents",
                    "bullets": [
                        "Dedicated agents for Requirements, Architecture, Code & QA.",
                        "Multi-agent collaboration via structured event orchestration.",
                        "Contextual awareness across the entire repository."
                    ]
                },
                {
                    "title": "Governance & Safety",
                    "metric": "100% Audit Control",
                    "bullets": [
                        "Human-in-the-loop review checkpoints at critical gates.",
                        "Full lineage tracking in PostgreSQL persistent database.",
                        "Zero unapproved modifications or security bypasses."
                    ]
                }
            ],
            "speaker_notes": "Welcome executive team. Today we present our Autonomous SDLC Platform, designed to transform enterprise software engineering by orchestrating 15 specialized AI agents under strict human governance."
        },
        {
            "category": "BUSINESS PROBLEM",
            "title": "Traditional SDLC Bottlenecks & Operational Friction",
            "subtitle": "High manual effort, inconsistent code quality, and delayed release cycles impact business agility.",
            "cards": [
                {
                    "title": "Slow Time-to-Market",
                    "metric": "6-8 Weeks Avg",
                    "bullets": [
                        "Manual requirement parsing and SRS document creation.",
                        "Siloed communication between Product, Engineering, and QA.",
                        "Backlog accumulation leading to delayed feature releases."
                    ]
                },
                {
                    "title": "Quality & Drift Risks",
                    "metric": "35% Rework Cost",
                    "bullets": [
                        "Requirements ambiguity leads to mid-sprint architecture changes.",
                        "Inconsistent test coverage across legacy components.",
                        "Security vulnerabilities identified late in staging pipelines."
                    ]
                },
                {
                    "title": "High Engineering Expense",
                    "metric": "60% Overhead",
                    "bullets": [
                        "Senior developers spend excess hours writing boilerplate.",
                        "Manual video walkthroughs and documentation maintenance.",
                        "Scaling engineering teams increases coordination overhead."
                    ]
                }
            ],
            "speaker_notes": "Enterprise teams face recurring bottlenecks: 6 to 8 week delivery cycles, 35% rework costs, and massive administrative overhead that slows down business innovation."
        },
        {
            "category": "PROPOSED SOLUTION",
            "title": "The EY Autonomous SDLC Platform Architecture",
            "subtitle": "An end-to-end multi-agent environment powering autonomous software development.",
            "cards": [
                {
                    "title": "Intelligent Intake & BRD",
                    "bullets": [
                        "Upload raw business notes or PDF documents.",
                        "AI Business Analyst agent generates verified BRD & SRS.",
                        "Instant gap analysis and risk identification."
                    ]
                },
                {
                    "title": "Autonomous Agent Workspace",
                    "bullets": [
                        "Agents generate production code, tests, and API schemas.",
                        "Real-time WebSocket progress tracking and logs.",
                        "Integrated diagram generation and visual inspection."
                    ]
                },
                {
                    "title": "Video & PPTX Generation",
                    "bullets": [
                        "Automatically renders HD executive presentation videos.",
                        "TTS voice-over narration with customizable presenter styles.",
                        "Direct export of PPTX decks and video artifacts."
                    ]
                }
            ],
            "speaker_notes": "Our platform solves this friction by combining document intake, multi-agent code generation, and automated HD video presentation generation into a unified workspace."
        },
        {
            "category": "BUSINESS BENEFITS & KPIS",
            "title": "Quantifiable Value & Return on Investment",
            "subtitle": "Accelerating delivery velocity while lowering operational risk and cost.",
            "cards": [
                {
                    "title": "Delivery Velocity",
                    "metric": "85% Faster",
                    "bullets": [
                        "Feature delivery timeline reduced from 6 weeks to 3 days.",
                        "Instant automated documentation and test generation.",
                        "Rapid iteration cycles with instant feedback loops."
                    ]
                },
                {
                    "title": "Quality & Security",
                    "metric": "99.4% Pass Rate",
                    "bullets": [
                        "Zero critical security defects reaching production.",
                        "100% test coverage across all generated endpoints.",
                        "Strict alignment with enterprise coding guidelines."
                    ]
                },
                {
                    "title": "Cost Efficiency",
                    "metric": "4.5x ROI",
                    "bullets": [
                        "Drastic reduction in developer manual overhead.",
                        "Higher engineering satisfaction and focus on innovation.",
                        "Estimated annual savings of $2.4M per engineering org."
                    ]
                }
            ],
            "speaker_notes": "In summary, the business impact is clear: 85% faster time-to-market, 99.4% security pass rates, and a demonstrated 4.5x return on investment."
        }
    ]

    # ── DECK 2: Technical Architecture ───────────────────────────────────────
    tech_slides = [
        {
            "category": "SYSTEM ARCHITECTURE",
            "title": "Enterprise Microservice & Agent Infrastructure",
            "subtitle": "Decoupled React SPA, FastAPI backend, PostgreSQL single source of truth, and local AI pipelines.",
            "cards": [
                {
                    "title": "Frontend Presentation",
                    "metric": "React 18 + Vite",
                    "bullets": [
                        "TypeScript SPA with modern dark glassmorphism styling.",
                        "Real-time WebSocket connection for live agent streaming.",
                        "Interactive slide canvas, video workspace, and script editor."
                    ]
                },
                {
                    "title": "Backend Agent Engine",
                    "metric": "FastAPI + Python 3.11",
                    "bullets": [
                        "Asynchronous Uvicorn server handling REST & WebSockets.",
                        "Autonomous multi-agent task runner and state management.",
                        "Local fallback capabilities with offline execution modes."
                    ]
                },
                {
                    "title": "Single Source of Truth",
                    "metric": "PostgreSQL Database",
                    "bullets": [
                        "ACID-compliant persistence for projects, slides & artifacts.",
                        "Versioned GeneratedArtifact storing exact JSON structures.",
                        "Full audit trail and execution history log."
                    ]
                }
            ],
            "speaker_notes": "The technical architecture is built for high reliability. A React 18 frontend communicates asynchronously with a FastAPI agent engine, backed by PostgreSQL as our single source of truth."
        },
        {
            "category": "MULTI-AGENT WORKFLOW",
            "title": "Autonomous Agent Collaboration & Event Pipeline",
            "subtitle": "Structured coordination between specialized agents to produce verified enterprise deliverables.",
            "cards": [
                {
                    "title": "1. Requirement & BA Agents",
                    "bullets": [
                        "Parses raw user prompts or uploaded PDF requirements.",
                        "Generates structured BRD, SRS, and feature backlogs.",
                        "Stores artifacts in PostgreSQL with strict schemas."
                    ]
                },
                {
                    "title": "2. Architecture & Code Agents",
                    "bullets": [
                        "Synthesizes system architecture and database models.",
                        "Generates production code, unit tests, and OpenAPI endpoints.",
                        "Validates code execution in isolated environments."
                    ]
                },
                {
                    "title": "3. Video & Media Pipeline",
                    "bullets": [
                        "Pillow & PPTX frame rendering for visual slides.",
                        "TTS audio synthesis with exact per-slide speaker_notes.",
                        "SadTalker / Cartoon presenter lip-sync compositing."
                    ]
                }
            ],
            "speaker_notes": "Our multi-agent workflow operates in three distinct phases: Requirement analysis, Code & Architecture generation, and Media presentation synthesis."
        },
        {
            "category": "TECHNOLOGY STACK",
            "title": "Modern Enterprise Tech Stack Components",
            "subtitle": "Proven open-source and enterprise-grade tools powering the application platform.",
            "cards": [
                {
                    "title": "Core Application Stack",
                    "bullets": [
                        "React 18, TypeScript, Tailwind CSS, Lucide Icons.",
                        "FastAPI, Python 3.11, Pydantic v2, SQLAlchemy ORM.",
                        "PostgreSQL 15 database engine with psycopg2."
                    ]
                },
                {
                    "title": "AI & Media Pipeline",
                    "bullets": [
                        "Google Antigravity & LLM APIs for intelligent agents.",
                        "Python-PPTX & Pillow for slide frame rasterization.",
                        "Edge-TTS & FFmpeg for audio/video compositing."
                    ]
                },
                {
                    "title": "DevOps & Local Execution",
                    "bullets": [
                        "Uvicorn async server on localhost port 8008.",
                        "Vite dev server on localhost port 5173.",
                        "Strict local git configuration and environment security."
                    ]
                }
            ],
            "speaker_notes": "Here is our complete technology stack: React, TypeScript, FastAPI, PostgreSQL, Python-PPTX, Edge-TTS, and FFmpeg."
        },
        {
            "category": "DEPLOYMENT & SECURITY",
            "title": "Enterprise Security, Isolation & Governance",
            "subtitle": "Ensuring strict data protection, untracked secret keys, and robust deployment compliance.",
            "cards": [
                {
                    "title": "Git & Secret Security",
                    "metric": "Zero Leak Policy",
                    "bullets": [
                        "Comprehensive .gitignore preventing key leakage.",
                        "Local environment untracking via git rm --cached.",
                        "Strict audit checks for all checked-in artifacts."
                    ]
                },
                {
                    "title": "Runtime Safety",
                    "metric": "Sandboxed Execution",
                    "bullets": [
                        "Isolated python environment for background agents.",
                        "Strict payload validation via Pydantic models.",
                        "Role-based access control and JWT authentication."
                    ]
                },
                {
                    "title": "Production Deployment",
                    "metric": "Cloud Native Ready",
                    "bullets": [
                        "Containerized with Docker & Kubernetes manifests.",
                        "Support for GCP Cloud Run and AWS ECS deployment.",
                        "Continuous integration & automated smoke tests."
                    ]
                }
            ],
            "speaker_notes": "Finally, deployment and security: we enforce zero secret leakage, sandboxed agent execution, and cloud-native container readiness."
        }
    ]

    path1 = create_deck("AI_SDLC_Executive_Overview.pptx", "AI SDLC Executive Overview", exec_slides)
    path2 = create_deck("AI_SDLC_Technical_Architecture.pptx", "AI SDLC Technical Architecture", tech_slides)

    print(f"DECK 1 SAVED: {path1}")
    print(f"DECK 2 SAVED: {path2}")

if __name__ == "__main__":
    generate_all_decks()
