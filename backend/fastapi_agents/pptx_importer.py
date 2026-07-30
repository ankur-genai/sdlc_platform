"""
pptx_importer.py
=================
Reverse of pptx_builder.py: reads an existing, user-supplied .pptx file back
into the editable slide-dict shape the workspace already understands (same
shape /video/render and /video/render/from-pdf return), so a user can upload
a deck they already have and keep editing/re-rendering it here instead of
starting from a blank AI-generated one.

Real-world .pptx files use semantic title/body placeholders, which
`slide.shapes.title` finds directly. But this app's own pptx_builder.py (and
many hand-built decks) use freeform, absolutely-positioned textboxes instead
of placeholders — there is no semantic "title" to query. For those, the
title is identified as the largest-font text shape in the upper portion of
the slide, and tiny shapes pinned to the bottom (footers, page numbers,
brand marks) are excluded from the body content entirely.
"""
from __future__ import annotations

import io
from .logging_config import get_logger
import re
from typing import Any

from pptx import Presentation

logger = get_logger(__name__)

# Shapes whose top edge falls in the bottom 8% of the slide are treated as
# footer/page-number/brand-mark noise, never real slide content.
_FOOTER_ZONE_FRACTION = 0.92

# Short strings that are near-universally decorative (brand marks, section
# eyebrow labels, confidentiality/date stamps) rather than real slide
# content, regardless of where they sit on the slide.
_NOISE_PATTERNS = (
    re.compile(r"^[A-Z]{1,4}$"),                       # bare brand initials, e.g. "EY"
    re.compile(r"^confidential\b", re.IGNORECASE),
    re.compile(r"^\d{1,3}\s*/\s*\d{1,3}$"),             # "02 / 12" page counters
)


def _is_noise(text: str) -> bool:
    t = text.strip()
    return bool(t) and any(p.match(t) for p in _NOISE_PATTERNS)


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()


def _max_font_size(shape) -> float:
    size = 0.0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                size = max(size, run.font.size.pt)
    return size


def _guess_layout(title_text: str, body_texts: list[str], slide_idx: int, total: int) -> str:
    """Best-effort mapping onto the editor's small layout vocabulary
    (title/content/two_column/quote/metric/closing)."""
    joined = " ".join(body_texts).strip()
    lower = (title_text + " " + joined).lower()
    if slide_idx == 0 and len(joined) < 40:
        return "title"
    if slide_idx == total - 1 and any(w in lower for w in ("thank you", "questions", "conclusion")):
        return "closing"
    if joined.strip().startswith(("“", '"', "'")) or "quote" in title_text.lower():
        return "quote"
    if len(body_texts) >= 2 and all(len(t) < 200 for t in body_texts):
        return "two_column"
    if not joined and title_text:
        return "title"
    return "content"


def _rasterize_pptx_slides(file_bytes: bytes) -> list[str]:
    """Convert each PPTX slide to a full-fidelity PNG data URI.

    Renderer priority (first success wins):
      1. PowerPoint COM  (Windows; comtypes/pywin32) — exact rendering, matches PowerPoint
      2. LibreOffice     (cross-platform; soffice)   — high-fidelity via headless PDF export
      3. Empty list      (no PIL placeholder; caller uses text content only)

    The PIL rectangle fallback has been intentionally removed: it produces
    fake coloured boxes that look nothing like the real slide, which is worse
    than showing no image at all.
    """
    import base64
    import shutil
    import tempfile
    import time
    from pathlib import Path

    images: list[str] = []

    # ── Method 1: PowerPoint COM via win32com (Windows, STA + message pump) ──
    # "Call was rejected by callee" (-2147418111) means PowerPoint's COM server
    # is busy. On an STA thread this happens when time.sleep() blocks the
    # Windows message pump — COM dispatch can never deliver the reply.
    # Fix: pump COM messages continuously while waiting.
    try:
        import win32com.client   # type: ignore  (pywin32)
        import pythoncom          # type: ignore
        import threading

        com_results: list[str] = []
        com_error:   list[str] = []

        def _pump_wait(seconds: float):
            """Sleep for `seconds` while pumping COM messages every 50 ms."""
            import time as _time
            end = _time.monotonic() + seconds
            while _time.monotonic() < end:
                pythoncom.PumpWaitingMessages()
                _time.sleep(0.05)

        def _com_export_worker():
            pythoncom.CoInitialize()          # STA mode for this thread
            tmpdir_w = tempfile.mkdtemp()
            try:
                pptx_path = Path(tmpdir_w) / "deck.pptx"
                pptx_path.write_bytes(file_bytes)

                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                try:
                    try:
                        powerpoint.DisplayAlerts = 1  # ppAlertsNone
                    except Exception:
                        pass

                    prs = powerpoint.Presentations.Open(str(pptx_path), -1, 0, 0)
                    slide_count = prs.Slides.Count
                    logger.info("[PptxImporter] COM/STA: opened %d slides", slide_count)

                    _pump_wait(0.5)

                    for i in range(1, slide_count + 1):
                        out_path = str(Path(tmpdir_w) / f"slide_{i:04d}.png")
                        max_attempts = 5
                        for attempt in range(max_attempts):
                            try:
                                pythoncom.PumpWaitingMessages()
                                prs.Slides(i).Export(out_path, "PNG", 1920, 1080)
                                break
                            except Exception as com_err:
                                if attempt < max_attempts - 1:
                                    wait = 0.5 * (2 ** attempt)
                                    logger.warning("[PptxImporter] COM slide %d attempt %d failed (%s) — retry in %.1fs",
                                                   i, attempt + 1, com_err, wait)
                                    _pump_wait(wait)
                                else:
                                    raise
                        png_bytes = Path(out_path).read_bytes()
                        b64 = base64.b64encode(png_bytes).decode("utf-8")
                        com_results.append(f"data:image/png;base64,{b64}")
                        logger.info("[PptxImporter] Renderer=PowerPoint_COM slide=%d size=%d bytes", i, len(png_bytes))
                        _pump_wait(0.1)

                    prs.Close()
                finally:
                    try:
                        powerpoint.Quit()
                    except Exception:
                        pass
                    _pump_wait(0.5)
            except Exception as exc:
                com_error.append(str(exc))
                logger.warning("[PptxImporter] COM/STA worker error: %s", exc)
            finally:
                shutil.rmtree(tmpdir_w, ignore_errors=True)
                pythoncom.CoUninitialize()

        t = threading.Thread(target=_com_export_worker, daemon=True)
        t.start()
        t.join(timeout=300)   # Allow up to 5 minutes for large decks

        if t.is_alive():
            logger.warning("[PptxImporter] PowerPoint COM timed out after 300s")
        elif com_error:
            logger.warning("[PptxImporter] PowerPoint COM failed: %s", com_error[0])
        elif com_results:
            logger.info("[PptxImporter] PowerPoint COM rasterization COMPLETE: %d slides", len(com_results))
            return com_results

    except ImportError as ie:
        logger.info("[PptxImporter] win32com/pywin32 not available (%s) — skipping PowerPoint COM", ie)
    except Exception as exc:
        logger.warning("[PptxImporter] PowerPoint COM rasterization failed: %s", exc)

    # ── Method 2: LibreOffice PPTX→PDF→PyMuPDF/pdf2image ─────────────────────
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if Path(candidate).exists():
            soffice = candidate
            break

    if soffice:
        import subprocess

        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "deck.pptx"
            pptx_path.write_bytes(file_bytes)
            try:
                subprocess.run(
                    [soffice, "--headless", "--norestore", "--convert-to", "pdf",
                     "--outdir", tmpdir, str(pptx_path)],
                    capture_output=True, timeout=120, check=True,
                )
                pdf_path = Path(tmpdir) / "deck.pdf"
                if pdf_path.exists():
                    # PyMuPDF path
                    try:
                        import fitz  # PyMuPDF
                        doc = fitz.open(str(pdf_path))
                        for page in doc:
                            pix = page.get_pixmap(dpi=150)
                            b64 = base64.b64encode(pix.tobytes("png")).decode("utf-8")
                            images.append(f"data:image/png;base64,{b64}")
                        if images:
                            logger.info("[PptxImporter] Renderer=LibreOffice+PyMuPDF  slides=%d", len(images))
                            return images
                    except Exception as exc:
                        logger.warning("[PptxImporter] PyMuPDF rasterization failed: %s", exc)

                    # pdf2image path
                    try:
                        import io
                        from pdf2image import convert_from_path
                        imgs = convert_from_path(str(pdf_path), dpi=150)
                        for img in imgs:
                            buf = io.BytesIO()
                            img.save(buf, format="PNG")
                            b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
                            images.append(f"data:image/png;base64,{b64}")
                        if images:
                            logger.info("[PptxImporter] Renderer=LibreOffice+pdf2image  slides=%d", len(images))
                            return images
                    except Exception as exc:
                        logger.warning("[PptxImporter] pdf2image rasterization failed: %s", exc)
            except Exception as exc:
                logger.warning("[PptxImporter] LibreOffice PPTX->PDF conversion failed: %s", exc)
    else:
        logger.warning("[PptxImporter] LibreOffice not found — PowerPoint COM was the only real renderer on this system")

    # ── No real renderer available ─────────────────────────────────────────────
    logger.error("[PptxImporter] ALL rasterizers failed. Returning no images (no PIL placeholder).")
    return []


def _render_pptx_slide_pil(slide, prs_width: int, prs_height: int, target_w: int = 1920, target_h: int = 1080) -> str:
    """Fallback PIL canvas renderer: renders PPTX slide shapes (boxes, cards, text, pictures, tables) to 1920x1080 PNG data URI."""
    import base64
    from PIL import Image, ImageDraw

    scale_x = target_w / (prs_width or 12192000)
    scale_y = target_h / (prs_height or 6858000)

    # Base canvas (dark theme default for EY presentation)
    img = Image.new("RGB", (target_w, target_h), color=(46, 46, 56)) # #2E2E38
    draw = ImageDraw.Draw(img)

    # Yellow left accent stripe
    draw.rectangle([0, 0, 16, target_h], fill=(255, 230, 0)) # #FFE600
    # Header & footer bars
    draw.rectangle([16, 0, target_w, 90], fill=(26, 26, 36)) # #1A1A24
    draw.rectangle([16, target_h - 52, target_w, target_h], fill=(26, 26, 36))

    for shape in slide.shapes:
        left = int((shape.left or 0) * scale_x)
        top = int((shape.top or 0) * scale_y)
        width = int((shape.width or 0) * scale_x)
        height = int((shape.height or 0) * scale_y)

        # Draw card background box for content shapes
        if width > 0 and height > 0 and top > 90 and top < target_h - 52:
            draw.rectangle([left, top, left + width, top + height], fill=(58, 58, 72), outline=(144, 144, 160))

        # Paste picture shapes if present
        if hasattr(shape, "image") and shape.image:
            try:
                pic_img = Image.open(io.BytesIO(shape.image.blob))
                pic_img = pic_img.resize((max(width, 10), max(height, 10)))
                img.paste(pic_img, (left, top))
            except Exception:
                pass

        # Draw text frame content
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()
            if text:
                text_color = (255, 255, 255) if top <= 90 else (216, 216, 224)
                draw.text((left + 12, top + 12), text[:300], fill=text_color)

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
    return f"data:image/png;base64,{b64}"


def parse_pptx_to_slides(file_bytes: bytes) -> list[dict[str, Any]]:
    """Parse an uploaded .pptx into the same slide-dict shape used elsewhere
    in the app: title, subtitle, content, speaker_notes, layout, duration, page_image."""
    import base64
    prs = Presentation(io.BytesIO(file_bytes))
    slide_width = prs.slide_width or 12192000
    slide_height = prs.slide_height or 6858000
    total = len(prs.slides)
    slides: list[dict[str, Any]] = []

    # Attempt full-deck WYSIWYG rasterization via LibreOffice + PyMuPDF/pdf2image first
    raster_images = _rasterize_pptx_slides(file_bytes)

    for idx, slide in enumerate(prs.slides):
        candidates: list[dict[str, Any]] = []
        extracted_diagram_image: str | None = None

        for shape in slide.shapes:
            # Check for embedded picture shapes / diagram artwork
            if not extracted_diagram_image:
                try:
                    if hasattr(shape, "image") and shape.image:
                        image_bytes = shape.image.blob
                        ext = shape.image.ext or "png"
                        b64 = base64.b64encode(image_bytes).decode("utf-8")
                        extracted_diagram_image = f"data:image/{ext};base64,{b64}"
                except Exception:
                    pass

            if not shape.has_text_frame:
                continue
            text = _shape_text(shape)
            if not text or _is_noise(text):
                continue
            top = shape.top if shape.top is not None else 0
            if top >= slide_height * _FOOTER_ZONE_FRACTION:
                continue  # footer / page-number / brand-mark strip — never content
            candidates.append({
                "text": text, "top": top, "font_size": _max_font_size(shape),
                "is_title_placeholder": shape == slide.shapes.title,
            })

        title_text = ""
        body_texts: list[str] = []
        if candidates:
            placeholder_title = next((c for c in candidates if c["is_title_placeholder"] and c["text"]), None)
            title_candidate = placeholder_title or max(candidates, key=lambda c: c["font_size"])
            title_text = title_candidate["text"].splitlines()[0][:120]
            remainder_lines = title_candidate["text"].splitlines()[1:]
            body_texts = (["\n".join(remainder_lines)] if remainder_lines else [])
            body_texts += [c["text"] for c in candidates if c is not title_candidate]
            # Keep original top-to-bottom reading order for the body.
            ordered = sorted(
                (c for c in candidates if c is not title_candidate),
                key=lambda c: c["top"],
            )
            body_texts = [c["text"] for c in ordered]
            if remainder_lines:
                body_texts = ["\n".join(remainder_lines)] + body_texts

        content = "\n".join(
            line if line.lstrip().startswith(("•", "-", "*")) else f"• {line}"
            for text in body_texts
            for line in text.splitlines()
            if line.strip()
        )

        speaker_notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            speaker_notes = slide.notes_slide.notes_text_frame.text.strip()

        layout = _guess_layout(title_text, body_texts, idx, total)

        page_image = None
        if idx < len(raster_images):
            page_image = raster_images[idx]
        # No PIL fallback: if the real rasterizer didn't produce an image for
        # this slide we leave page_image as None.  The frontend will show the
        # text-only content layout instead of fake placeholder rectangles.

        slide_dict: dict[str, Any] = {
            "title": title_text or f"Slide {idx + 1}",
            "subtitle": "",
            "content": content,
            "speaker_notes": speaker_notes,
            "layout": layout,
            "duration": 25,
        }
        if page_image:
            slide_dict["page_image"] = page_image
        if extracted_diagram_image and "page_image" not in slide_dict:
            slide_dict["diagram_image"] = extracted_diagram_image

        has_p = bool(slide_dict.get("page_image"))
        p_len = len(slide_dict.get("page_image") or "")
        has_d = bool(slide_dict.get("diagram_image"))
        d_len = len(slide_dict.get("diagram_image") or "")
        logger.info("[DEBUG_BACKEND] PPTX Slide %d: title=%r, page_image_exists=%s, len=%d, diagram_image_exists=%s, len=%d",
                    idx + 1, slide_dict.get("title"), has_p, p_len, has_d, d_len)

        slides.append(slide_dict)

    if not slides:
        raise ValueError("No slides found in the uploaded .pptx file")

    return slides
