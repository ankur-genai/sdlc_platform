"""
pptx_progressive.py
====================
Builds progressive "reveal stage" variants of a user-*imported* .pptx —
duplicating each slide in place (same theme/images/layout, nothing
rebuilt) and truncating its content shapes so the deck can be rasterized
once and played back as a build-in animation, without ever re-theming the
user's own artwork (see video_pipeline_local.py's PptxFrameRenderer, which
this is designed to feed).

Unlike a deck generated from scratch by this platform (where each slide is
a `dict` of structured fields — title/content/table/chart — that can be
truncated by simply passing less data to the theme renderer), an imported
.pptx's "bullets" are just arbitrary shapes at arbitrary positions: some
decks put a whole bullet list in one placeholder's paragraphs, others (like
pptxgenjs-authored decks) give every bullet/label its own single-paragraph
text box. `_group_revealable_shapes` handles both by clustering shapes
into rows using their vertical position rather than assuming a paragraph
structure — inherently a heuristic, not true semantic understanding of the
deck's structure, so the reveal grouping approximates "top-to-bottom in
waves" rather than the author's exact intended bullet groupings.
"""
from __future__ import annotations

import copy
import io
from typing import Any, List, Optional, Tuple

from .logging_config import get_logger

logger = get_logger(__name__)

_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
_R_ATTRS = ("embed", "id", "link", "cs", "dm", "lo", "qs")

# Cap on stages per slide — matches the synthetic-deck path's cap (see
# video_pipeline_local.py's _reveal_stage_count): more than a handful of
# build-in beats stops feeling "subtle, business-friendly" per the product
# requirement and starts feeling like a slow striptease of bullet points.
MAX_STAGES_PER_SLIDE = 4


def _is_title_placeholder(shape) -> bool:
    try:
        if not shape.is_placeholder:
            return False
        from pptx.enum.shapes import PP_PLACEHOLDER
        return shape.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE,
        )
    except Exception:
        return False


def _shape_center(shape, slide_w: int, slide_h: int) -> Tuple[float, float]:
    return (
        ((shape.left or 0) + (shape.width or 0) / 2) / slide_w,
        ((shape.top or 0) + (shape.height or 0) / 2) / slide_h,
    )


def _group_revealable_shapes(slide, title_text: str, subtitle_text: str) -> List[List[Any]]:
    """Semantic-first reveal grouping. Primarily follows the slide's own
    shape *document order* (each shape's position in `<p:spTree>`, which is
    `slide.shapes`'s natural iteration order) rather than raw Y-position —
    for a two-column slide (e.g. body paragraphs on the left, a stat/
    achievement panel on the right), Y-position clustering interleaves the
    two columns whenever a left-column and right-column shape happen to
    sit at a similar height, scrambling the reveal order relative to how
    the deck's own narration is almost always written (one section fully,
    then the next) — document order doesn't have that problem, since a
    deck's author (or a generator script) places a section's shapes
    together in the underlying XML.

    Two real semantic signals are honored when present:
    - a top-level *group* shape (`<p:grpSp>`, an explicit "these belong
      together" from the author) is kept as a single, whole reveal unit —
      not split into its children;
    - a shape with a TITLE/CENTER_TITLE/SUBTITLE placeholder type is always
      excluded from the reveal (shown from stage 0), same as the header
      text is today.

    Most decks (anything built by a generation script like pptxgenjs, which
    is what this heuristic was actually developed against) have neither —
    plain, ungrouped text boxes — so document order plus a position-based
    adjacency check (did two document-order-consecutive shapes land near
    each other, or is this a jump to a new section?) is the fallback and
    does most of the real work. `title_text`/`subtitle_text` are accepted
    but unused — the platform's own text extraction can mis-identify which
    shape is the "title" (see video_pipeline_local.py's PPTX import path),
    so a placeholder-type check (or, lacking one, absolute position) is the
    only reliable signal."""
    slide_w = getattr(slide.part.package.presentation_part.presentation, "slide_width", 12192000)
    slide_h = getattr(slide.part.package.presentation_part.presentation, "slide_height", 6858000)
    header_band = slide_h * 0.15   # ~15% from the top: title/subtitle/slide-number badge
    footer_band = slide_h * 0.90   # ~10% from the bottom: page-number/branding footers

    candidates: List[Any] = []
    for shape in slide.shapes:
        if _is_title_placeholder(shape):
            continue
        has_text = shape.has_text_frame and shape.text_frame.text.strip()
        # Non-text shapes (pictures, charts, tables, SmartArt graphic
        # frames) are still real reveal-worthy content — a slide's image
        # shouldn't just always be visible while its caption builds in.
        is_visual = getattr(shape, "shape_type", None) is not None and not has_text
        if not has_text and not is_visual:
            continue
        top = shape.top if shape.top is not None else 0
        if top <= header_band or top >= footer_band:
            continue
        candidates.append(shape)

    if len(candidates) <= 1:
        return []

    # Adjacency threshold: two document-order-consecutive shapes join the
    # same group when their centers are this close (as a fraction of the
    # slide's diagonal) — close enough to read as "one visual unit" (an
    # icon next to its label) rather than "the next distinct section".
    adjacency_threshold = 0.18

    groups: List[List[Any]] = [[candidates[0]]]
    prev_center = _shape_center(candidates[0], slide_w, slide_h)
    for shape in candidates[1:]:
        center = _shape_center(shape, slide_w, slide_h)
        dist = ((center[0] - prev_center[0]) ** 2 + (center[1] - prev_center[1]) ** 2) ** 0.5
        if dist > adjacency_threshold:
            groups.append([])
        groups[-1].append(shape)
        prev_center = center

    # Cap at MAX_STAGES_PER_SLIDE reveal groups (not counting the always-
    # visible header) by merging the smallest adjacent groups first — this
    # merges within the now-correct document/reading order, so a cap-driven
    # merge still can't accidentally fuse two unrelated sections the way a
    # position-only approach could.
    while len(groups) > MAX_STAGES_PER_SLIDE:
        merge_at = min(range(len(groups) - 1), key=lambda i: len(groups[i]) + len(groups[i + 1]))
        groups[merge_at] = groups[merge_at] + groups[merge_at + 1]
        del groups[merge_at + 1]

    return groups


def _duplicate_slide(prs, source_slide):
    """Deep-copies every shape from `source_slide` into a newly added slide
    (same layout), remapping relationship ids so embedded images/media keep
    resolving. This is the one part of python-pptx that has no built-in
    support (see the pptx skill's own gotchas) — the two failure modes to
    guard against are (1) leftover placeholder shapes the new slide inherits
    from its layout, and (2) dangling r:embed/r:id references if a shape's
    relationship isn't also copied onto the new slide's own rels part."""
    from pptx.oxml.ns import qn

    dest = prs.slides.add_slide(source_slide.slide_layout)
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    rid_map = {}
    for rId, rel in source_slide.part.rels.items():
        try:
            if rel.is_external:
                new_rid = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
            else:
                new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rId] = new_rid
        except Exception:
            continue

    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        for el in new_el.iter():
            for local in _R_ATTRS:
                attr = f"{_R_NS}{local}"
                old_rid = el.get(attr)
                if old_rid and old_rid in rid_map:
                    el.set(attr, rid_map[old_rid])
        dest.shapes._spTree.append(new_el)

    src_cSld = source_slide._element.find(qn("p:cSld"))
    src_bg = src_cSld.find(qn("p:bg")) if src_cSld is not None else None
    if src_bg is not None:
        dest_cSld = dest._element.find(qn("p:cSld"))
        dest_cSld.insert(0, copy.deepcopy(src_bg))

    return dest


def _sldId_for_slide(prs, slide):
    from pptx.oxml.ns import qn
    partname = slide.part.partname
    for sldId in prs.slides._sldIdLst:
        rId = sldId.get(qn("r:id"))
        rel = prs.part.rels[rId]
        if rel.target_part.partname == partname:
            return sldId
    raise LookupError("slide not found in presentation's sldIdLst")


def _remove_shape_text(shape) -> None:
    """Blanks a shape's own text without touching its paragraph/run
    formatting — used for the pre-reveal stages of any shape not yet in
    scope for the current stage."""
    for p in list(shape.text_frame.paragraphs):
        p._p.getparent().remove(p._p)
    # A text body needs at least one (even if empty) paragraph to stay valid.
    from pptx.oxml.ns import qn
    txBody = shape.text_frame._txBody
    if txBody.find(qn("a:p")) is None:
        from pptx.oxml import parse_xml
        empty_p = parse_xml(
            '<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
        txBody.append(empty_p)


def _group_geometry(group: List[Any], slide_w: int, slide_h: int) -> Optional[dict]:
    """Bounding-box geometry of `group`'s shapes, as fractions of the slide
    (0..1): `cx`/`cy` (centroid, for camera panning — see _focus_pan_exprs)
    and `x0`/`y0`/`x1`/`y1` (union bbox, for sizing the zoom to actually
    frame the region — see _focus_zoom_target — and for the spotlight/dim
    effect on everything outside it). None if the group is empty or the
    slide has no size info."""
    if not group or not slide_w or not slide_h:
        return None
    x0 = y0 = float("inf")
    x1 = y1 = float("-inf")
    for shape in group:
        try:
            left, top = shape.left or 0, shape.top or 0
            right, bottom = left + (shape.width or 0), top + (shape.height or 0)
        except Exception:
            continue
        x0, y0 = min(x0, left), min(y0, top)
        x1, y1 = max(x1, right), max(y1, bottom)
    if x0 == float("inf"):
        return None
    return {
        "cx": (x0 + x1) / 2 / slide_w, "cy": (y0 + y1) / 2 / slide_h,
        "x0": x0 / slide_w, "y0": y0 / slide_h,
        "x1": x1 / slide_w, "y1": y1 / slide_h,
    }


def build_progressive_import_deck(original_bytes: bytes,
                                   slides_meta: List[dict]) -> Optional[Tuple[bytes, List[int], List[List[Optional[dict]]], List[List[str]]]]:
    """Given the raw bytes of a user-imported .pptx and the platform's own
    already-parsed metadata for each slide (title/subtitle — kept in the
    signature but unused, see _group_revealable_shapes), returns
    `(expanded_pptx_bytes, stage_counts, focus_geometry, stage_texts)`:
    - `stage_counts[i]` is how many stage-slides original slide `i` expanded
      into (1 = untouched, in its original position). The caller rasterizes
      the expanded deck in one LibreOffice pass and slices the resulting
      frames back per original slide using `stage_counts` — mirroring the
      from-scratch synthetic-deck path in video_pipeline_local.py so the two
      share the same "one soffice conversion per video" performance profile.
    - `focus_geometry[i]` is a list (len stage_counts[i]) of geometry dicts
      (see _group_geometry: cx/cy centroid + x0/y0/x1/y1 bbox, all as slide
      fractions) or `None` — whatever content is newly visible at that
      stage, so the video's camera can pan/zoom toward it and a spotlight
      can dim everything else (see FFmpegComposer._make_progressive_clip
      and _apply_spotlight in video_pipeline_local.py) instead of a
      generic pan.
    - `stage_texts[i]` is a list (len stage_counts[i]) of each stage's own
      on-screen text, so narration sentences can be matched to whichever
      stage they actually describe (_best_narration_partition in
      video_pipeline_local.py) instead of a blind even split across stages
      — the reveal grouping being in the right *order* (see
      _group_revealable_shapes) doesn't by itself guarantee a 1:1
      correspondence between narration sentence count and reveal-group
      count, so content matching still matters even with correct ordering.

    Returns None if nothing in the deck has more than one reveal row (so
    the caller keeps its existing single-frame-per-slide behavior), or on
    any failure — this must never risk corrupting a render that already
    works today by half-editing the user's real file.
    """
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(original_bytes))
    except Exception as exc:
        logger.warning("[PptxProgressive] Could not open imported .pptx: %s", exc)
        return None

    original_count = len(prs.slides)
    if original_count != len(slides_meta):
        logger.warning("[PptxProgressive] slide count mismatch (%d pptx vs %d metadata); skipping",
                       original_count, len(slides_meta))
        return None

    groups_per_slide: List[List[List[Any]]] = []
    for slide, meta in zip(prs.slides, slides_meta):
        try:
            groups = _group_revealable_shapes(slide, meta.get("title", ""), meta.get("subtitle", ""))
        except Exception as exc:
            logger.debug("[PptxProgressive] Grouping failed for a slide (%s); leaving it static", exc)
            groups = []
        groups_per_slide.append(groups)

    stage_counts = [len(g) + 1 if len(g) > 1 else 1 for g in groups_per_slide]
    if all(sc == 1 for sc in stage_counts):
        return None

    slide_w, slide_h = prs.slide_width, prs.slide_height
    focus_points: List[List[Optional[dict]]] = []
    stage_texts: List[List[str]] = []
    for groups, sc in zip(groups_per_slide, stage_counts):
        if sc <= 1:
            focus_points.append([None])
            stage_texts.append([""])
            continue
        # Duplicate stage k's newly-visible content is groups[k]; the final
        # (untouched) stage matches the last duplicate's content exactly
        # (see the stage/group mapping below), so it reuses that geometry.
        points = [_group_geometry(groups[k], slide_w, slide_h) for k in range(sc - 1)]
        points.append(points[-1] if points else None)
        focus_points.append(points)
        # The group's own on-screen text — lets the caller match narration
        # sentences to whichever stage's content they actually describe
        # (see _best_narration_partition in video_pipeline_local.py) instead
        # of blindly splitting narration into even chunks by stage index.
        texts = [
            " ".join(s.text_frame.text.strip() for s in groups[k] if s.has_text_frame and s.text_frame.text.strip())
            for k in range(sc - 1)
        ]
        texts.append(texts[-1] if texts else "")
        stage_texts.append(texts)

    try:
        sldIdLst = prs.slides._sldIdLst
        original_sldId_elements = list(sldIdLst)

        # Walk slides back-to-front: inserting/reordering earlier slides
        # doesn't shift the positions of slides we haven't processed yet.
        for slide_idx in range(original_count - 1, -1, -1):
            groups = groups_per_slide[slide_idx]
            sc = stage_counts[slide_idx]
            if sc <= 1:
                continue

            source_slide = prs.slides[slide_idx]
            new_slides = [_duplicate_slide(prs, source_slide) for _ in range(sc - 1)]

            anchor = original_sldId_elements[slide_idx]
            for dup in new_slides:
                dup_sldId = _sldId_for_slide(prs, dup)
                sldIdLst.remove(dup_sldId)
                anchor.addprevious(dup_sldId)

            # Stage k (0-based, k < sc-1) shows groups[0..k]; the untouched
            # original slide (stage sc-1) is the full, unedited reveal.
            for stage_idx, dup in enumerate(new_slides):
                shapes_to_blank = [shp for grp in groups[stage_idx + 1:] for shp in grp]
                dup_shape_by_name = {s.name: s for s in dup.shapes if s.has_text_frame}
                for shp in shapes_to_blank:
                    dup_shape = dup_shape_by_name.get(shp.name)
                    if dup_shape is not None:
                        try:
                            _remove_shape_text(dup_shape)
                        except Exception as exc:
                            logger.debug("[PptxProgressive] Could not blank a shape (%s)", exc)

        out = io.BytesIO()
        prs.save(out)
        return out.getvalue(), stage_counts, focus_points, stage_texts
    except Exception as exc:
        logger.warning("[PptxProgressive] Expansion failed (%s); rendering single frame per slide", exc)
        return None
